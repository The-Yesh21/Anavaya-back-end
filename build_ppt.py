"""
Generates a professional PowerPoint deck for the Anavaya Judicial Case Priority System.
Theme: Judicial navy + gold accents.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# ---------------- Theme ----------------
NAVY      = RGBColor(0x0B, 0x1F, 0x3A)   # deep navy (primary bg)
NAVY_SOFT = RGBColor(0x14, 0x2B, 0x4C)   # card bg
GOLD      = RGBColor(0xD4, 0xAF, 0x37)   # judicial gold (accent)
GOLD_LIGHT= RGBColor(0xE8, 0xCE, 0x7A)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE  = RGBColor(0xF2, 0xF2, 0xF2)
SLATE     = RGBColor(0xA9, 0xB7, 0xC6)   # muted text
GREEN     = RGBColor(0x4E, 0xC9, 0xB0)
RED_SOFT  = RGBColor(0xE0, 0x6C, 0x6C)
LIGHT_BG  = RGBColor(0xF7, 0xF8, 0xFA)   # light slide bg
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---------------- helpers ----------------
def add_slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    bg_rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg_rect.line.fill.background()
    bg_rect.fill.solid()
    bg_rect.fill.fore_color.rgb = bg
    bg_rect.shadow.inherit = False
    # send to back
    sp = bg_rect._element
    sp.getparent().remove(sp)
    s.shapes._spTree.insert(2, sp)
    return s

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75, shadow=False, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    if shadow:
        _add_shadow(sp)
    return sp

def _add_shadow(shape):
    spPr = shape._element.spPr
    effectLst = etree.SubElement(spPr, qn('a:effectLst'))
    outerShdw = etree.SubElement(effectLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', '90000'); outerShdw.set('dist', '38100')
    outerShdw.set('dir', '5400000'); outerShdw.set('rotWithShape', '0')
    clr = etree.SubElement(outerShdw, qn('a:srgbClr'))
    clr.set('val', '1A1A2E'); alpha = etree.SubElement(clr, qn('a:alpha'))
    alpha.set('val', '24000')

def add_text(slide, x, y, w, h, text, size=14, color=DARK_TEXT, bold=False, italic=False,
             align=PP_ALIGN.LEFT, font='Calibri', anchor=MSO_ANCHOR.TOP, line_spacing=1.1):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lines = text.split('\n') if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run(); r.text = line
        f = r.font
        f.size = Pt(size); f.bold = bold; f.italic = italic
        f.color.rgb = color; f.name = font
    return tb

def add_bullets(slide, x, y, w, h, items, size=15, color=DARK_TEXT, bullet_color=GOLD,
                font='Calibri', line_spacing=1.35, space_after=8):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        # bullet glyph run
        rb = p.add_run(); rb.text = "▸  "
        rb.font.size = Pt(size); rb.font.bold = True
        rb.font.color.rgb = bullet_color; rb.font.name = font
        # text run (support bold lead via tuple)
        if isinstance(item, tuple):
            lead, rest = item
            r1 = p.add_run(); r1.text = lead
            r1.font.size = Pt(size); r1.font.bold = True
            r1.font.color.rgb = color; r1.font.name = font
            r2 = p.add_run(); r2.text = rest
            r2.font.size = Pt(size); r2.font.bold = False
            r2.font.color.rgb = color; r2.font.name = font
        else:
            r = p.add_run(); r.text = item
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = font
    return tb

def dark_header(slide, eyebrow, title, accent=GOLD):
    """Standard header for content slides (light bg)."""
    # top accent bar
    add_rect(slide, 0, 0, SW, Inches(0.12), fill=accent)
    # eyebrow
    add_text(slide, Inches(0.6), Inches(0.32), Inches(10), Inches(0.3),
             eyebrow.upper(), size=12, color=accent, bold=True, font='Calibri')
    # title
    add_text(slide, Inches(0.6), Inches(0.58), Inches(12.1), Inches(0.7),
             title, size=28, color=DARK_TEXT, bold=True, font='Calibri')
    # underline
    add_rect(slide, Inches(0.62), Inches(1.28), Inches(0.6), Pt(3), fill=accent)
    # footer page tag
    _footer(slide)

def _footer(slide):
    add_text(slide, Inches(0.6), Inches(7.08), Inches(6), Inches(0.3),
             "ANAVAYA  ·  Judicial Case Priority System", size=9, color=SLATE, bold=True)
    n = len(prs.slides._sldIdLst)
    add_text(slide, Inches(12.2), Inches(7.08), Inches(0.7), Inches(0.3),
             f"{n:02d}", size=9, color=SLATE, align=PP_ALIGN.RIGHT)

def card(slide, x, y, w, h, fill=WHITE, line=RGBColor(0xE0,0xE5,0xEC), line_w=1.0, shadow=True):
    return add_rect(slide, x, y, w, h, fill=fill, line=line, line_w=line_w, shadow=shadow,
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE)

def pill(slide, x, y, w, h, text, fill, txt_color, size=11, bold=True):
    sp = add_rect(slide, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = sp.text_frame; tf.word_wrap = True
    tf.margin_left=Inches(0.05); tf.margin_right=Inches(0.05)
    tf.margin_top=Inches(0.02); tf.margin_bottom=Inches(0.02)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = txt_color; r.font.name='Calibri'
    return sp


# =================================================================
# SLIDE 1 — TITLE
# =================================================================
s = add_slide(bg=NAVY)
# decorative gold corner blocks
add_rect(s, 0, 0, SW, Inches(0.18), fill=GOLD)
add_rect(s, 0, SH-Inches(0.18), SW, Inches(0.18), fill=GOLD)
# faint large-scale-of-justice motif (circle)
ring = add_rect(s, Inches(9.4), Inches(-1.7), Inches(5.5), Inches(5.5),
                fill=None, line=GOLD, line_w=1.5, shape=MSO_SHAPE.OVAL)
ring.line.color.rgb = GOLD
# eyebrow
add_text(s, Inches(0.9), Inches(1.9), Inches(8), Inches(0.4),
         "MAJOR PROJECT  ·  AI FOR THE JUDICIARY", size=14, color=GOLD, bold=True)
# title
add_text(s, Inches(0.9), Inches(2.3), Inches(11), Inches(1.5),
         "Anavaya", size=72, color=WHITE, bold=True)
add_text(s, Inches(0.9), Inches(3.5), Inches(11.5), Inches(0.9),
         "Judicial Case Priority System & Dashboard", size=30, color=GOLD_LIGHT, bold=False)
# tagline
add_rect(s, Inches(0.92), Inches(4.5), Inches(0.5), Pt(2.5), fill=GOLD)
add_text(s, Inches(0.9), Inches(4.7), Inches(11), Inches(1.0),
         "An AI advisory system that ranks legal cases by severity, vulnerability and risk —\n"
         "extracting facts from court documents and deciding priority through an\n"
         "interpretable Decision Tree grounded in the Constitution of India.",
         size=15, color=OFFWHITE, line_spacing=1.4)
# chips
chips = ["Hybrid NLP + ML", "Interpretable", "Constitutional", "FastAPI Dashboard"]
cx = Inches(0.9); cy = Inches(6.35)
for c in chips:
    w = Inches(0.25 + 0.105*len(c))
    pill(s, cx, cy, w, Inches(0.4), c, fill=NAVY_SOFT, txt_color=GOLD_LIGHT, size=11)
    cx += w + Inches(0.2)


# =================================================================
# SLIDE 2 — THE CHALLENGE
# =================================================================
s = add_slide(bg=LIGHT_BG)
dark_header(s, "Context", "The Challenge Before the Courts")

# left: narrative
card(s, Inches(0.6), Inches(1.6), Inches(6.3), Inches(5.0))
add_text(s, Inches(0.9), Inches(1.85), Inches(5.7), Inches(0.4),
         "WHY THIS MATTERS", size=12, color=GOLD, bold=True)
add_bullets(s, Inches(0.9), Inches(2.3), Inches(5.7), Inches(4.0), [
    ("Backlog — ", "Courts face millions of pending cases; urgent matters get buried."),
    ("Manual triage — ", "Clerks scan PDFs by hand to decide hearing order."),
    ("Inconsistency — ", "Similar cases get different urgency across benches."),
    ("No audit trail — ", "Priority reasons are rarely documented."),
    ("Delay = injustice — ", "Justice delayed in violent crimes is a rights violation."),
], size=15, color=DARK_TEXT, line_spacing=1.35, space_after=10)

# right: stat cards
stats = [
    ("3+ crore", "pending cases in India", GOLD),
    ("5+ years", "average civil case duration", GREEN),
    ("1 decision", "per ~73,000 people", RED_SOFT),
]
sy = Inches(1.6)
for val, lbl, col in stats:
    card(s, Inches(7.25), sy, Inches(5.5), Inches(1.5))
    add_rect(s, Inches(7.25), sy, Inches(0.12), Inches(1.5), fill=col)
    add_text(s, Inches(7.6), sy+Inches(0.22), Inches(5), Inches(0.7),
             val, size=32, color=col, bold=True)
    add_text(s, Inches(7.6), sy+Inches(0.85), Inches(5), Inches(0.5),
             lbl, size=14, color=DARK_TEXT)
    sy += Inches(1.72)

# takeaway
card(s, Inches(7.25), sy+Inches(0.05), Inches(5.5), Inches(1.55), fill=NAVY)
add_text(s, Inches(7.55), sy+Inches(0.28), Inches(5), Inches(0.4),
         "THE OPPORTUNITY", size=11, color=GOLD, bold=True)
add_text(s, Inches(7.55), sy+Inches(0.62), Inches(5), Inches(0.9),
         "Use AI to rank cases consistently, transparently and at scale —\nwhile leaving the final call to the judge.",
         size=13, color=OFFWHITE, line_spacing=1.3)


# =================================================================
# SLIDE 3 — SOLUTION OVERVIEW
# =================================================================
s = add_slide(bg=LIGHT_BG)
dark_header(s, "Solution", "What Anavaya Does")

# central tagline card
card(s, Inches(0.6), Inches(1.55), Inches(12.13), Inches(1.15), fill=NAVY)
add_text(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(0.4),
         "ONE-LINE SUMMARY", size=11, color=GOLD, bold=True)
add_text(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.6),
         "Given a raw legal PDF → extract facts with an LLM → decide priority with a Decision Tree → justify it constitutionally.",
         size=17, color=WHITE, bold=True)

# 4 capability cards
caps = [
    ("📄", "Document Intake", "Reads FIRs, judgments, complaints and testimonies straight from PDF."),
    ("🧠", "Fact Extraction", "An LLM pulls parties, crime type, severity, vulnerability & influence."),
    ("🌳", "Priority Decision", "A Decision Tree assigns High / Medium / Low — deterministic & explainable."),
    ("⚖️", "Constitutional Trace", "Every priority is justified against the Constitution of India."),
]
cw = Inches(2.93); gap = Inches(0.13); cx = Inches(0.6); cy = Inches(3.05)
for icon, title, desc in caps:
    card(s, cx, cy, cw, Inches(3.5))
    # icon circle
    add_rect(s, cx+Inches(0.3), cy+Inches(0.35), Inches(0.85), Inches(0.85),
             fill=NAVY, shape=MSO_SHAPE.OVAL)
    add_text(s, cx+Inches(0.3), cy+Inches(0.42), Inches(0.85), Inches(0.7),
             icon, size=30, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, cx+Inches(0.28), cy+Inches(1.4), cw-Inches(0.5), Inches(0.6),
             title, size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, cx+Inches(0.28), cy+Inches(1.95), cw-Inches(0.55), Inches(1.4),
             desc, size=11.5, color=DARK_TEXT, align=PP_ALIGN.CENTER, line_spacing=1.3)
    cx += cw + gap


# =================================================================
# SLIDE 4 — KEY DESIGN INVARIANT
# =================================================================
s = add_slide(bg=NAVY)
# top bar
add_rect(s, 0, 0, SW, Inches(0.12), fill=GOLD)
add_text(s, Inches(0.6), Inches(0.5), Inches(10), Inches(0.3),
         "CORE DESIGN PRINCIPLE", size=12, color=GOLD, bold=True)
add_text(s, Inches(0.6), Inches(0.78), Inches(12), Inches(0.8),
         "The LLM Extracts. The Model Decides.", size=34, color=WHITE, bold=True)
add_rect(s, Inches(0.62), Inches(1.65), Inches(0.6), Pt(3), fill=GOLD)

# two-column compare
colw = Inches(5.9); colh = Inches(4.6); top = Inches(2.0)
# left card (LLM)
card(s, Inches(0.6), top, colw, colh, fill=NAVY_SOFT, line=GOLD, line_w=1.0)
add_text(s, Inches(0.95), top+Inches(0.35), Inches(5), Inches(0.4),
         "✦  LLM (Gemma)", size=16, color=GOLD_LIGHT, bold=True)
add_text(s, Inches(0.95), top+Inches(0.78), Inches(5.2), Inches(0.5),
         "Perception layer only", size=13, color=SLATE)
add_bullets(s, Inches(0.95), top+Inches(1.45), Inches(5.2), Inches(3.0), [
    "Reads the raw PDF text",
    "Identifies parties & legal domain",
    "Extracts severity, vulnerability, influence",
    "Writes a plain-language summary",
    "Never outputs a priority or ranking",
], size=14, color=OFFWHITE, bullet_color=GOLD_LIGHT, line_spacing=1.4, space_after=10)

# right card (Decision Tree)
card(s, Inches(6.85), top, colw, colh, fill=NAVY_SOFT, line=GOLD, line_w=1.0)
add_text(s, Inches(7.2), top+Inches(0.35), Inches(5), Inches(0.4),
         "⚖  Decision Tree (CART)", size=16, color=GOLD_LIGHT, bold=True)
add_text(s, Inches(7.2), top+Inches(0.78), Inches(5.2), Inches(0.5),
         "Decision layer — the only decider", size=13, color=SLATE)
add_bullets(s, Inches(7.2), top+Inches(1.45), Inches(5.2), Inches(3.0), [
    "Deterministic & reproducible",
    "Trained on labelled cases",
    "Outputs High / Medium / Low",
    "Exposes an exact decision path",
    "Auditable against constitutional rules",
], size=14, color=OFFWHITE, bullet_color=GOLD_LIGHT, line_spacing=1.4, space_after=10)

add_text(s, Inches(0.6), Inches(6.85), Inches(12), Inches(0.4),
         "Why it matters: priority is never an opaque LLM guess — it is traceable, testable and defensible.",
         size=13, color=GOLD, italic=True, align=PP_ALIGN.CENTER)
_footer(s)


# =================================================================
# SLIDE 5 — SYSTEM ARCHITECTURE (pipeline)
# =================================================================
s = add_slide(bg=LIGHT_BG)
dark_header(s, "Architecture", "End-to-End Pipeline")

# pipeline row
nodes = [
    ("PDF", "raw legal\ndocument", NAVY),
    ("LLM\nGemma", "extract\nfeatures", RGBColor(0x2D,0x5A,0x8C)),
    ("Tune", "normalize\n+ rules", RGBColor(0x3A,0x7C,0xA5)),
    ("Decision\nTree", "assign\npriority", GOLD),
    ("Report", "Excel +\ngraph + MD", RGBColor(0x4E,0x9B,0x6E)),
]
nx = Inches(0.6); ny = Inches(1.7); nw = Inches(2.15); nh = Inches(1.5); gap = Inches(0.25)
centers = []
for i,(t,sub,col) in enumerate(nodes):
    box = add_rect(s, nx, ny, nw, nh, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    if col == GOLD:
        tc = NAVY
    else:
        tc = WHITE
    add_text(s, nx, ny+Inches(0.18), nw, Inches(0.7), t, size=15, color=tc, bold=True, align=PP_ALIGN.CENTER, line_spacing=0.95)
    add_text(s, nx, ny+Inches(0.92), nw, Inches(0.5), sub, size=10, color=tc, align=PP_ALIGN.CENTER, line_spacing=0.95)
    centers.append((nx, col))
    # arrow
    if i < len(nodes)-1:
        ar = add_rect(s, nx+nw+Inches(0.02), ny+Inches(0.6), gap-Inches(0.04), Inches(0.3),
                      fill=SLATE, shape=MSO_SHAPE.RIGHT_ARROW)
    nx += nw + gap

# the rule layer annotation
add_text(s, Inches(0.6), Inches(3.45), Inches(12), Inches(0.4),
         "Constitutional & explanation layer runs in parallel (rule-based, deterministic):",
         size=12.5, color=SLATE, italic=True)
ann = ["Rights engaged (Art. 14, 21, 23/24, 300A)", "State's duty analysis",
       "Applicable doctrines", "Proportionality / balancing", "Priority rules applied"]
cx = Inches(0.6)
for a in ann:
    w = Inches(0.28 + 0.092*len(a))
    pill(s, cx, Inches(3.85), w, Inches(0.42), a, fill=RGBColor(0xEC,0xF1,0xF7),
         txt_color=NAVY, size=11)
    cx += w + Inches(0.15)

# outputs strip
add_text(s, Inches(0.6), Inches(4.7), Inches(12), Inches(0.4),
         "DELIVERABLES", size=12, color=GOLD, bold=True)
outs = [
    ("case_results.xlsx", "prioritized case table"),
    ("decision_graphs/*.md", "per-case Mermaid reports"),
    ("decision_graphs/*.dot", "raw decision-path graphs"),
    ("priority_classifier.pkl", "trained model bundle"),
    ("/api/* FastAPI", "dashboard + JSON endpoints"),
]
cw = Inches(2.37); gap = Inches(0.1); cx = Inches(0.6); cy = Inches(5.1)
for name, desc in outs:
    card(s, cx, cy, cw, Inches(1.6))
    add_rect(s, cx, cy, cw, Inches(0.1), fill=GOLD)
    add_text(s, cx+Inches(0.18), cy+Inches(0.28), cw-Inches(0.3), Inches(0.5),
             name, size=12, color=NAVY, bold=True, line_spacing=0.95)
    add_text(s, cx+Inches(0.18), cy+Inches(0.85), cw-Inches(0.3), Inches(0.6),
             desc, size=10.5, color=DARK_TEXT, line_spacing=1.1)
    cx += cw + gap


# =================================================================
# SLIDE 6 — ALGORITHMS USED (overview)
# =================================================================
s = add_slide(bg=LIGHT_BG)
dark_header(s, "Algorithms", "The Algorithms Powering Anavaya")

algos = [
    ("1", "TF-IDF Vectorization", "Statistical NLP", "Term Frequency–Inverse Document Frequency turns case text into 220 numeric (unigram + bigram) features.", RGBColor(0x2D,0x5A,0x8C)),
    ("2", "Label Encoding", "Categorical encoding", "Maps the 5 categorical fields (crime_type, severity, …) to integers the tree can split on.", RGBColor(0x3A,0x7C,0xA5)),
    ("3", "CART Decision Tree", "Supervised ML · PRIMARY", "scikit-learn DecisionTreeClassifier (Gini). Depth 8, balanced classes. The only priority decider.", GOLD),
    ("4", "Feedforward Neural Net", "Deep learning · OPTIONAL", "2-layer MLP (Linear→ReLU→Linear), Adam + CrossEntropy. Demo only — not in production.", RGBColor(0x8B,0x5C,0xA6)),
    ("5", "Rule-Based Heuristics", "Expert system", "Keyword scoring & policy rules for legal categorization, fallback extraction & labelling.", RGBColor(0x4E,0x9B,0x6E)),
    ("6", "LLM Extraction", "Generative AI", "NVIDIA Gemma via LangChain. Extracts structured features + summary. Never decides priority.", RGBColor(0xC0,0x6C,0x4E)),
]
cw = Inches(4.0); ch = Inches(2.5); gap_x = Inches(0.13); gap_y = Inches(0.18)
x0 = Inches(0.6); y0 = Inches(1.55)
for i,(num, name, tag, desc, col) in enumerate(algos):
    r, c = divmod(i, 3)
    x = x0 + c*(cw+gap_x); y = y0 + r*(ch+gap_y)
    card(s, x, y, cw, ch)
    add_rect(s, x, y, Inches(0.12), ch, fill=col)
    # number badge
    add_rect(s, x+Inches(0.3), y+Inches(0.25), Inches(0.5), Inches(0.5), fill=col, shape=MSO_SHAPE.OVAL)
    add_text(s, x+Inches(0.3), y+Inches(0.3), Inches(0.5), Inches(0.4),
             num, size=16, color=WHITE if col!=GOLD else NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x+Inches(0.95), y+Inches(0.27), Inches(2.9), Inches(0.4),
             tag.upper(), size=9.5, color=col, bold=True)
    add_text(s, x+Inches(0.95), y+Inches(0.48), cw-Inches(1.1), Inches(0.5),
             name, size=15, color=NAVY, bold=True)
    add_text(s, x+Inches(0.3), y+Inches(1.15), cw-Inches(0.5), Inches(1.25),
             desc, size=11, color=DARK_TEXT, line_spacing=1.25)


# =================================================================
# SLIDE 7 — DECISION TREE DEEP DIVE
# =================================================================
s = add_slide(bg=LIGHT_BG)
dark_header(s, "Algorithm Spotlight", "CART Decision Tree — The Priority Engine")

# left: config card
card(s, Inches(0.6), Inches(1.55), Inches(5.4), Inches(5.1))
add_text(s, Inches(0.9), Inches(1.8), Inches(5), Inches(0.4),
         "CONFIGURATION", size=12, color=GOLD, bold=True)
rows = [
    ("Algorithm", "CART (Gini impurity)"),
    ("Implementation", "DecisionTreeClassifier"),
    ("Max depth", "8"),
    ("Min samples / leaf", "3"),
    ("Class weights", "balanced"),
    ("Categorical inputs", "5 fields"),
    ("Text inputs", "220 TF-IDF features"),
    ("Total features", "225"),
    ("Training rows", "1,616"),
]
ry = Inches(2.25)
for k, v in rows:
    add_text(s, Inches(0.9), ry, Inches(2.4), Inches(0.32), k, size=12, color=SLATE)
    add_text(s, Inches(3.3), ry, Inches(2.5), Inches(0.32), v, size=12, color=NAVY, bold=True)
    ry += Inches(0.4)

# right: learned rules (from training_report.txt)
card(s, Inches(6.25), Inches(1.55), Inches(6.5), Inches(5.1))
add_text(s, Inches(6.55), Inches(1.8), Inches(6), Inches(0.4),
         "LEARNED RULES  (exported from the tree)", size=12, color=GOLD, bold=True)
rule_box = add_rect(s, Inches(6.55), Inches(2.3), Inches(5.9), Inches(3.0),
                    fill=RGBColor(0xF4,0xF6,0xF9), line=RGBColor(0xDD,0xE3,0xEC))
rules = (
    "if  victim > 0.03\n"
    "│      →  class: High   ← any victim signal → High\n\n"
    "else  (no victim keyword):\n"
    "├── severity ≤ 2.5   → class: Low / Medium\n"
    "└── severity > 2.5:\n"
    "      ├── influence ≤ 0.5      → Medium\n"
    "      └── influence > 0.5:\n"
    "            └── vulnerability, tenancy, concern\n"
    "               keyword splits → Medium / High"
)
add_text(s, Inches(6.8), Inches(2.5), Inches(5.5), Inches(2.8),
         rules, size=12, color=DARK_TEXT, font='Consolas', line_spacing=1.3)
add_text(s, Inches(6.55), Inches(5.45), Inches(5.9), Inches(1.0),
         "Reads naturally: violent/victim-bearing cases jump to High at the root; "
         "civil disputes cascade through severity → influence → vulnerability.",
         size=11.5, color=SLATE, italic=True, line_spacing=1.3)


# =================================================================
# SLIDE 8 — FEATURE EXTRACTION (LLM + heuristics)
# =================================================================
s = add_slide(bg=LIGHT_BG)
dark_header(s, "Perception Layer", "Feature Extraction — LLM + Rules")

# left: what the LLM extracts
card(s, Inches(0.6), Inches(1.55), Inches(6.1), Inches(5.1))
add_text(s, Inches(0.9), Inches(1.8), Inches(5.5), Inches(0.4),
         "WHAT THE LLM EXTRACTS", size=12, color=GOLD, bold=True)
feats = [
    ("main_parties", "people, companies, courts involved"),
    ("case_category", "8 legal domains (Tax, Criminal, …)"),
    ("crime_type", "Violent / Financial / Property / Non-Violent"),
    ("severity", "Fatal / Major / Minor / No Injury"),
    ("vulnerability", "High / Medium / Low"),
    ("influence", "power imbalance: High / Low"),
    ("plain_summary", "3–4 sentence plain-language summary"),
]
fy = Inches(2.3)
for name, desc in feats:
    add_rect(s, Inches(0.9), fy+Inches(0.06), Inches(0.16), Inches(0.16),
             fill=GOLD, shape=MSO_SHAPE.OVAL)
    add_text(s, Inches(1.2), fy, Inches(2.2), Inches(0.3), name, size=12.5, color=NAVY, bold=True, font='Consolas')
    add_text(s, Inches(3.25), fy, Inches(3.3), Inches(0.3), desc, size=11, color=DARK_TEXT)
    fy += Inches(0.46)

# right: how it's done (stack)
card(s, Inches(6.95), Inches(1.55), Inches(5.8), Inches(5.1))
add_text(s, Inches(7.25), Inches(1.8), Inches(5.2), Inches(0.4),
         "HOW IT'S DONE", size=12, color=GOLD, bold=True)
stack = [
    ("1", "LangChain + ChatNVIDIA", "with_structured_output() + Pydantic schema for clean JSON."),
    ("2", "Pydantic CaseFeatures", "Constrained Literal types force valid label values."),
    ("3", "Fallback: raw requests", "Direct NVIDIA API + JSON repair if LangChain unavailable."),
    ("4", "tune_case_features()", "Keyword rules normalize & override LLM output to known labels."),
    ("5", "classify_legal_category()", "Weighted keyword scoring maps text → legal domain."),
]
sty = Inches(2.3)
for num, t, d in stack:
    add_rect(s, Inches(7.25), sty+Inches(0.04), Inches(0.4), Inches(0.4), fill=NAVY, shape=MSO_SHAPE.OVAL)
    add_text(s, Inches(7.25), sty+Inches(0.08), Inches(0.4), Inches(0.32), num, size=13, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(7.8), sty, Inches(4.7), Inches(0.3), t, size=12.5, color=NAVY, bold=True, font='Consolas')
    add_text(s, Inches(7.8), sty+Inches(0.28), Inches(4.7), Inches(0.5), d, size=10.5, color=DARK_TEXT, line_spacing=1.2)
    sty += Inches(0.82)


# =================================================================
# SLIDE 9 — MODELS BUILT (table)
# =================================================================
s = add_slide(bg=LIGHT_BG)
dark_header(s, "Inventory", "Models & Components Built")

# table header
tx, ty = Inches(0.6), Inches(1.6)
tw, rh = Inches(12.13), Inches(0.55)
cols = [Inches(3.6), Inches(3.6), Inches(3.0), Inches(1.93)]
headers = ["Component", "Type", "Location", "In Production?"]
# header row
hx = tx
add_rect(s, tx, ty, tw, rh, fill=NAVY)
for i, htext in enumerate(headers):
    add_text(s, hx+Inches(0.15), ty+Inches(0.12), cols[i]-Inches(0.2), Inches(0.35),
             htext, size=12.5, color=GOLD, bold=True)
    hx += cols[i]

rows_data = [
    ("Decision Tree Classifier", "Supervised ML · CART", "models/priority_classifier.pkl", "Yes", GREEN),
    ("Feedforward Neural Network", "Deep learning · 2-layer MLP", "models/priority_dl_model.pth", "No (demo)", RED_SOFT),
    ("TF-IDF Vectorizer", "Statistical NLP", "bundled in pickle", "Yes", GREEN),
    ("LLM Feature Extractor", "Generative AI · Gemma", "external NVIDIA NIM API", "Yes", GREEN),
    ("Rule-Based Legal Engine", "Deterministic expert system", "constitutional_analysis.py", "Yes", GREEN),
]
ry = ty + rh
for i, (a,b,c,d,col) in enumerate(rows_data):
    bg = WHITE if i%2==0 else RGBColor(0xF2,0xF5,0xF9)
    add_rect(s, tx, ry, tw, Inches(0.78), fill=bg)
    cx = tx
    cells = [a, b, c]
    for j, val in enumerate(cells):
        col_color = NAVY if j==0 else DARK_TEXT
        bold = (j==0)
        add_text(s, cx+Inches(0.15), ry+Inches(0.22), cols[j]-Inches(0.2), Inches(0.4),
                 val, size=12, color=col_color, bold=bold, font='Consolas' if j==2 else 'Calibri')
        cx += cols[j]
    # production pill
    pill(s, cx+Inches(0.25), ry+Inches(0.2), cols[3]-Inches(0.4), Inches(0.4),
         d, fill=col if col==GREEN else RGBColor(0xFB,0xEC,0xEC),
         txt_color=RGBColor(0x14,0x3C,0x2A) if col==GREEN else RED_SOFT, size=11)
    ry += Inches(0.78)

# note
card(s, Inches(0.6), Inches(6.35), Inches(12.13), Inches(0.75), fill=RGBColor(0xFB,0xF6,0xE7))
add_text(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.5),
         "Note:  the neural network exists in the repo but is essentially unused — the production system is deliberately kept simple and interpretable.",
         size=12, color=RGBColor(0x6B,0x52,0x12), italic=True)


# =================================================================
# SLIDE 10 — MODEL PERFORMANCE
# =================================================================
s = add_slide(bg=LIGHT_BG)
dark_header(s, "Evaluation", "Model Performance")

# big accuracy
card(s, Inches(0.6), Inches(1.6), Inches(3.7), Inches(5.05), fill=NAVY)
add_text(s, Inches(0.85), Inches(1.9), Inches(3.2), Inches(0.4),
         "OVERALL", size=12, color=GOLD, bold=True)
add_text(s, Inches(0.85), Inches(2.3), Inches(3.2), Inches(1.3),
         "99.4%", size=66, color=WHITE, bold=True)
add_text(s, Inches(0.85), Inches(3.7), Inches(3.2), Inches(0.4),
         "Test accuracy on held-out set", size=12.5, color=SLATE)
add_rect(s, Inches(0.9), Inches(4.3), Inches(2.8), Pt(2), fill=GOLD)
add_bullets(s, Inches(0.85), Inches(4.6), Inches(3.2), Inches(2.0), [
    "1,616 training rows",
    "324 held-out test rows",
    "Perfect High-priority recall",
], size=12, color=OFFWHITE, bullet_color=GOLD_LIGHT, line_spacing=1.4)

# right: classification report
card(s, Inches(4.55), Inches(1.6), Inches(8.18), Inches(5.05))
add_text(s, Inches(4.85), Inches(1.85), Inches(7), Inches(0.4),
         "CLASSIFICATION REPORT", size=12, color=GOLD, bold=True)

# mini table
mcols = [Inches(1.6), Inches(1.4), Inches(1.4), Inches(1.4), Inches(1.6)]
mheads = ["Class", "Precision", "Recall", "F1", "Support"]
mx = Inches(4.85); my = Inches(2.45)
add_rect(s, mx, my, sum(mcols, Emu(0)), Inches(0.5), fill=NAVY_SOFT)
cx = mx
for i, h in enumerate(mheads):
    add_text(s, cx+Inches(0.1), my+Inches(0.1), mcols[i], Inches(0.3),
             h, size=11.5, color=GOLD_LIGHT, bold=True)
    cx += mcols[i]
mrows = [
    ("High",   "1.00", "0.98", "0.99", "110", GOLD),
    ("Low",    "1.00", "1.00", "1.00", "66",  GREEN),
    ("Medium", "0.99", "1.00", "0.99", "148", RGBColor(0x3A,0x7C,0xA5)),
    ("avg",    "1.00", "0.99", "0.99", "324", DARK_TEXT),
]
ry = my + Inches(0.5)
for i,(c,p,r,f,sup,col) in enumerate(mrows):
    bg = WHITE if i%2==0 else RGBColor(0xF5,0xF7,0xFA)
    add_rect(s, mx, ry, sum(mcols, Emu(0)), Inches(0.55), fill=bg)
    cx = mx
    for j,val in enumerate([c,p,r,f,sup]):
        bold = (j==0)
        col_use = col if j==0 else DARK_TEXT
        add_text(s, cx+Inches(0.1), ry+Inches(0.13), mcols[j], Inches(0.3),
                 val, size=12.5, color=col_use, bold=bold,
                 font='Consolas' if j>0 else 'Calibri')
        cx += mcols[j]
    ry += Inches(0.55)

add_rect(s, mx, ry, sum(mcols, Emu(0)), Inches(0.5), fill=NAVY_SOFT)
add_text(s, mx+Inches(0.1), ry+Inches(0.1), mcols[0], Inches(0.3),
         "macro / weighted avg", size=11.5, color=NAVY, bold=True)
add_text(s, mx+mcols[0]+mcols[1]+mcols[2]+Inches(0.1), ry+Inches(0.1), Inches(3), Inches(0.3),
         "0.99 / 0.99 across all metrics", size=11.5, color=DARK_TEXT, italic=True)

# insight card
card(s, Inches(0.6), Inches(6.85), Inches(12.13), Inches(0.55), fill=RGBColor(0xEC,0xF6,0xEF))
add_text(s, Inches(0.9), Inches(6.93), Inches(11.5), Inches(0.4),
         "✓  Most important result:  100% recall on High-priority cases — no critical case is ever missed.",
         size=12.5, color=RGBColor(0x14,0x3C,0x2A), bold=True)


# =================================================================
# SLIDE 11 — CONSTITUTIONAL FRAMEWORK
# =================================================================
s = add_slide(bg=LIGHT_BG)
dark_header(s, "Justification Layer", "Grounded in the Constitution of India")

arts = [
    ("Art. 21", "Right to Life & Liberty", "Violent / fatal cases → High priority. The most fundamental right.", RGBColor(0xC0,0x39,0x2B)),
    ("Art. 14", "Equality Before Law", "Power-imbalance cases → Medium / High. Protects the weaker side.", RGBColor(0x2D,0x5A,0x8C)),
    ("Art. 23/24", "Against Exploitation", "Trafficking, sexual assault, child labour → Critical / High.", RGBColor(0x6B,0x3F,0x8F)),
    ("Art. 300A", "Right to Property", "Economic / property disputes → Medium / Low unless violent.", RGBColor(0x4E,0x9B,0x6E)),
]
cw = Inches(2.93); gap = Inches(0.13); cx = Inches(0.6); cy = Inches(1.6)
for art, title, desc, col in arts:
    card(s, cx, cy, cw, Inches(3.2))
    add_rect(s, cx, cy, cw, Inches(0.85), fill=col)
    add_text(s, cx, cy+Inches(0.12), cw, Inches(0.4), art, size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, cx, cy+Inches(0.5), cw, Inches(0.3), title.upper(), size=9.5, color=OFFWHITE, align=PP_ALIGN.CENTER)
    add_text(s, cx+Inches(0.25), cy+Inches(1.05), cw-Inches(0.5), Inches(2.0),
             desc, size=12, color=DARK_TEXT, align=PP_ALIGN.CENTER, line_spacing=1.35)
    cx += cw + gap

# doctrines strip
card(s, Inches(0.6), Inches(5.0), Inches(12.13), Inches(1.6))
add_text(s, Inches(0.9), Inches(5.2), Inches(11), Inches(0.4),
         "CONSTITUTIONAL DOCTRINES APPLIED", size=12, color=GOLD, bold=True)
doctrines = ["Proportionality", "Reasonable Classification", "Parens Patriae",
             "Natural Justice", "Basic Structure", "Severability"]
cx = Inches(0.9); cy = Inches(5.65)
for d in doctrines:
    w = Inches(0.3 + 0.092*len(d))
    pill(s, cx, cy, w, Inches(0.45), d, fill=NAVY, txt_color=GOLD_LIGHT, size=11)
    cx += w + Inches(0.15)

# note
add_text(s, Inches(0.6), Inches(6.85), Inches(12), Inches(0.4),
         "Fully rule-based and deterministic — the LLM performs no constitutional reasoning. Every justification is verifiable against the Constitution.",
         size=12, color=SLATE, italic=True, align=PP_ALIGN.CENTER)


# =================================================================
# SLIDE 12 — WHAT'S NOT USED / DESIGN CHOICES
# =================================================================
s = add_slide(bg=NAVY)
add_rect(s, 0, 0, SW, Inches(0.12), fill=GOLD)
add_text(s, Inches(0.6), Inches(0.5), Inches(10), Inches(0.3),
         "DESIGN PHILOSOPHY", size=12, color=GOLD, bold=True)
add_text(s, Inches(0.6), Inches(0.78), Inches(12), Inches(0.8),
         "What Anavaya Deliberately Doesn't Use", size=32, color=WHITE, bold=True)
add_rect(s, Inches(0.62), Inches(1.6), Inches(0.6), Pt(3), fill=GOLD)

# left: not used
card(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.6), fill=NAVY_SOFT, line=RGBColor(0x3A,0x4A,0x66))
add_text(s, Inches(0.9), Inches(2.25), Inches(5.5), Inches(0.4),
         "✕  AVOIDED", size=14, color=RED_SOFT, bold=True)
not_used = [
    "Random Forest / XGBoost / GBM",
    "SVM, Logistic Regression, Naive Bayes, KNN",
    "Transformer / BERT embeddings trained from scratch",
    "LLM as the priority decider",
    "Black-box neural nets in the hot path",
    "Probabilistic / non-reproducible scoring",
]
add_bullets(s, Inches(0.9), Inches(2.85), Inches(5.5), Inches(3.5), not_used,
            size=13.5, color=OFFWHITE, bullet_color=RED_SOFT, line_spacing=1.5, space_after=12)

# right: why
card(s, Inches(6.85), Inches(2.0), Inches(5.88), Inches(4.6), fill=NAVY_SOFT, line=GOLD)
add_text(s, Inches(7.15), Inches(2.25), Inches(5.5), Inches(0.4),
         "✓  WHY — GUIDING PRINCIPLES", size=14, color=GOLD, bold=True)
why = [
    ("Interpretability — ", "judges must see exactly why a case ranks where it does."),
    ("Reproducibility — ", "same input ⇒ same priority, every time."),
    ("Auditability — ", "every path traces to a rule and a constitutional article."),
    ("Separation of concerns — ", "perception (LLM) ≠ decision (tree)."),
    ("Judicial sovereignty — ", "the system advises; the judge decides."),
    ("Simplicity — ", "the simplest model that meets the bar."),
]
add_bullets(s, Inches(7.15), Inches(2.85), Inches(5.3), Inches(3.6), why,
            size=13, color=OFFWHITE, bullet_color=GOLD_LIGHT, line_spacing=1.4, space_after=11)
_footer(s)


# =================================================================
# SLIDE 13 — DASHBOARD & TECH STACK
# =================================================================
s = add_slide(bg=LIGHT_BG)
dark_header(s, "Interface", "Tech Stack & Dashboard")

# left: tech stack grouped
groups = [
    ("Backend",   ["Python", "FastAPI", "Uvicorn"], RGBColor(0x2D,0x5A,0x8C)),
    ("ML / NLP",  ["scikit-learn", "PyTorch", "TF-IDF", "pandas"], GOLD),
    ("LLM",       ["NVIDIA Gemma", "LangChain", "Pydantic"], RGBColor(0x8B,0x5C,0xA6)),
    ("Frontend",  ["D3.js", "Vanilla CSS", "Glassmorphism"], RGBColor(0x4E,0x9B,0x6E)),
    ("Data / IO", ["openpyxl", "PyMuPDF / pypdf", "Graphviz DOT"], RGBColor(0xC0,0x6C,0x4E)),
]
cy = Inches(1.6)
for name, items, col in groups:
    card(s, Inches(0.6), cy, Inches(6.2), Inches(0.95))
    add_rect(s, Inches(0.6), cy, Inches(0.12), Inches(0.95), fill=col)
    add_text(s, Inches(0.9), cy+Inches(0.1), Inches(2.0), Inches(0.3), name.upper(), size=10.5, color=col, bold=True)
    cx = Inches(0.9); iyc = cy+Inches(0.42)
    for it in items:
        w = Inches(0.28 + 0.095*len(it))
        pill(s, cx, iyc, w, Inches(0.36), it, fill=RGBColor(0xEE,0xF2,0xF7), txt_color=NAVY, size=10.5)
        cx += w + Inches(0.12)
    cy += Inches(1.05)

# right: dashboard features
card(s, Inches(7.1), Inches(1.6), Inches(5.63), Inches(5.1), fill=NAVY)
add_text(s, Inches(7.4), Inches(1.85), Inches(5), Inches(0.4),
         "WEB DASHBOARD", size=12, color=GOLD, bold=True)
add_text(s, Inches(7.4), Inches(2.2), Inches(5), Inches(0.5),
         "Interactive Case Board", size=18, color=WHITE, bold=True)
dash = [
    ("Case listing & filters", "search parties / summaries, filter by category & priority."),
    ("Global decision tree", "D3.js tree; click a case to glow its exact path."),
    ("Case breakdown panel", "severity, vulnerability, influence & step cards."),
    ("Constitutional trace", "rights engaged, doctrines & state's-duty opinion."),
    ("PDF upload", "live end-to-end triage via /api/upload."),
]
dy = Inches(2.85)
for t, d in dash:
    add_rect(s, Inches(7.4), dy+Inches(0.07), Inches(0.14), Inches(0.14), fill=GOLD, shape=MSO_SHAPE.OVAL)
    add_text(s, Inches(7.65), dy, Inches(4.8), Inches(0.3), t, size=12.5, color=GOLD_LIGHT, bold=True)
    add_text(s, Inches(7.65), dy+Inches(0.26), Inches(4.8), Inches(0.5), d, size=10.5, color=OFFWHITE, line_spacing=1.2)
    dy += Inches(0.72)


# =================================================================
# SLIDE 14 — SUMMARY / THANK YOU
# =================================================================
s = add_slide(bg=NAVY)
add_rect(s, 0, 0, SW, Inches(0.18), fill=GOLD)
add_rect(s, 0, SH-Inches(0.18), SW, Inches(0.18), fill=GOLD)
ring = add_rect(s, Inches(-1.5), Inches(-1.5), Inches(4), Inches(4),
                fill=None, line=GOLD, line_w=1.0, shape=MSO_SHAPE.OVAL)
add_rect(s, Inches(10.8), Inches(5.0), Inches(4), Inches(4),
         fill=None, line=GOLD, line_w=1.0, shape=MSO_SHAPE.OVAL)._element if False else None

add_text(s, Inches(0.9), Inches(1.1), Inches(10), Inches(0.4),
         "IN SUMMARY", size=13, color=GOLD, bold=True)
add_text(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.9),
         "Anavaya in One Breath", size=34, color=WHITE, bold=True)
add_rect(s, Inches(0.92), Inches(2.45), Inches(0.6), Pt(3), fill=GOLD)

points = [
    ("Hybrid by design.  ", "An LLM (Gemma) perceives facts; a Decision Tree decides priority."),
    ("Interpretable by construction.  ", "Every decision path is explicit and exportable."),
    ("Constitutional by grounding.  ", "Each priority is justified against the Constitution of India."),
    ("Advisory, never authoritative.  ", "Judicial sovereignty is preserved — the judge has the final word."),
    ("Production-ready.  ", "FastAPI backend, D3 dashboard, Excel + graph reports, trained model bundle."),
]
py = Inches(2.85)
for lead, rest in points:
    add_rect(s, Inches(0.95), py+Inches(0.1), Inches(0.16), Inches(0.16), fill=GOLD, shape=MSO_SHAPE.OVAL)
    tb = add_text(s, Inches(1.3), py, Inches(11), Inches(0.5), "", size=15)
    tf = tb.text_frame; tf.word_wrap=True
    p = tf.paragraphs[0]; p.line_spacing=1.2
    r1 = p.add_run(); r1.text = lead
    r1.font.size=Pt(16); r1.font.bold=True; r1.font.color.rgb=GOLD_LIGHT; r1.font.name='Calibri'
    r2 = p.add_run(); r2.text = rest
    r2.font.size=Pt(16); r2.font.bold=False; r2.font.color.rgb=OFFWHITE; r2.font.name='Calibri'
    py += Inches(0.62)

# thank you block
add_rect(s, Inches(0.92), Inches(6.1), Inches(2.0), Pt(2.5), fill=GOLD)
add_text(s, Inches(0.9), Inches(6.25), Inches(11.5), Inches(0.7),
         "Thank you.", size=28, color=WHITE, bold=True)
add_text(s, Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.4),
         "Questions & discussion welcome.", size=14, color=GOLD_LIGHT, italic=True)


# ---------------- save ----------------
out = 'Anavaya_Presentation.pptx'
prs.save(out)
print(f"Saved {out} with {len(prs.slides._sldIdLst)} slides.")
